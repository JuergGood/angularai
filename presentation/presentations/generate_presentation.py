import os
import re
from typing import Any, Dict, List, Optional, Tuple, Union

import yaml
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER, PP_PLACEHOLDER_TYPE
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt


SlideData = Dict[str, Any]


# ----------------------------
# Path + media helpers
# ----------------------------

def _abs_from_plan(plan_path: str, maybe_rel: str) -> str:
    """Resolve path relative to the markdown plan file."""
    if not maybe_rel:
        return maybe_rel
    if os.path.isabs(maybe_rel):
        return maybe_rel
    return os.path.normpath(os.path.join(os.path.dirname(plan_path), maybe_rel))


def _fit_image_into_box(img_path: str, box_w_emu: int, box_h_emu: int) -> Tuple[int, int]:
    """Return (w_emu, h_emu) preserving aspect ratio to fit within the box (contain)."""
    with Image.open(img_path) as im:
        w_px, h_px = im.size
    if w_px <= 0 or h_px <= 0:
        return box_w_emu, box_h_emu

    img_ar = w_px / h_px
    box_ar = box_w_emu / box_h_emu

    if img_ar >= box_ar:
        w = box_w_emu
        h = int(box_w_emu / img_ar)
    else:
        h = box_h_emu
        w = int(box_h_emu * img_ar)
    return w, h


# ----------------------------
# YAML robustness helpers
# ----------------------------

_TOP_LEVEL_KEY_RE = re.compile(r"^[A-Za-z_]\w*:\s*(.*)$")
_BLOCK_SCALAR_START_RE = re.compile(r"^[A-Za-z_]\w*:\s*\|\s*$")


def _auto_indent_block_scalars(yaml_body: str) -> str:
    """
    Makes YAML more forgiving for 'key: |' blocks by ensuring subsequent lines
    are indented at least 2 spaces until the next top-level key.

    Fixes common mistakes like:
      content: |
      - item
    to:
      content: |
        - item

    Note: This does not fix structurally incorrect indentation where a new key
    (e.g. 'right:') was accidentally written inside a block scalar. For that,
    see _dedent_accidental_keys_in_block_scalars().
    """
    lines = yaml_body.splitlines()
    out: List[str] = []
    in_block = False

    for line in lines:
        if _BLOCK_SCALAR_START_RE.match(line):
            in_block = True
            out.append(line)
            continue

        # end block when we see a new top-level key (no leading spaces)
        if in_block and line and not line.startswith((" ", "\t")) and _TOP_LEVEL_KEY_RE.match(line):
            in_block = False

        if in_block:
            if line.strip() and not line.startswith((" ", "\t")):
                out.append("  " + line)
            else:
                out.append(line)
        else:
            out.append(line)

    return "\n".join(out)


def _dedent_accidental_keys_in_block_scalars(yaml_body: str, keys: List[str] = ["right:", "left:", "image:", "caption:", "text:", "mode:", "layout:"]) -> str:
    """
    Best-effort fix for mistakes like:

      left: |
        - bullet
        right:
          image: x.png

    where 'right:' was intended as a top-level key but is indented within the left block scalar.

    Strategy:
      - Track when we're inside a block scalar.
      - If we encounter a line with two+ leading spaces that (after strip) begins with one of the known keys,
        we dedent it to column 0 and also dedent subsequent more-indented lines (the nested mapping) by 2 spaces
        until we hit a non-indented top-level key.

    This is conservative and only targets known keys.
    """
    key_prefixes = tuple(keys)
    lines = yaml_body.splitlines()
    out: List[str] = []
    in_block = False
    pending_dedent = False

    for i, line in enumerate(lines):
        if _BLOCK_SCALAR_START_RE.match(line):
            in_block = True
            pending_dedent = False
            out.append(line)
            continue

        # If we hit a new top-level key, we are no longer in the block scalar
        if in_block and line and not line.startswith((" ", "\t")) and _TOP_LEVEL_KEY_RE.match(line):
            in_block = False
            pending_dedent = False
            out.append(line)
            continue

        if in_block:
            stripped = line.lstrip()
            indent = len(line) - len(stripped)

            if indent >= 2 and stripped.startswith(key_prefixes):
                # Dedent this line to become a top-level key
                out.append(stripped)
                pending_dedent = True
                continue

            if pending_dedent:
                # Dedent nested lines by 2 spaces if possible
                if line.startswith("  "):
                    out.append(line[2:])
                    continue
                else:
                    # If we hit a non-indented line inside pending_dedent, stop dedenting
                    pending_dedent = False
                    out.append(line)
                    continue

            out.append(line)
        else:
            out.append(line)

    return "\n".join(out)


# ----------------------------
# PPTX layout + placeholder helpers
# ----------------------------

def _norm_layout_name(name: str) -> str:
    return re.sub(r"\s+", " ", (name or "").strip().lower())


def _get_layout_index_by_name(prs: Presentation, wanted: str) -> Optional[int]:
    wn = _norm_layout_name(wanted)
    for i, layout in enumerate(prs.slide_layouts):
        if _norm_layout_name(layout.name) == wn:
            return i
    return None


def _get_layout_index_by_any_name(prs: Presentation, wanted_names: List[str]) -> Optional[int]:
    for n in wanted_names:
        idx = _get_layout_index_by_name(prs, n)
        if idx is not None:
            return idx
    return None


def _find_placeholder(slide, idx: Optional[int] = None, placeholder_type: Optional[int] = None):
    """Best-effort placeholder lookup by index or placeholder type."""
    if idx is not None:
        try:
            return slide.placeholders[idx]
        except Exception:
            return None
    if placeholder_type is not None:
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.type == placeholder_type:
                    return ph
            except Exception:
                continue
    return None


def _remove_shape(shape) -> None:
    try:
        el = shape.element
        el.getparent().remove(el)
    except Exception:
        pass


def _remove_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides.sldId_lst)
    xml_slides.remove(slides[index])

# ----------------------------
# Generator
# ----------------------------

class PresentationGenerator:
    """
    Key properties:
      - Uses layouts ONLY from template.pptx by layout name (no index fallbacks).
      - Parses each slide block as YAML.
      - Adds layout: agenda.
      - Unifies all two-content variants into:
            layout: two_content
            mode: image_left | image_right | text_both
        Additionally supports legacy aliases:
            layout: image_left / image_right / text_both
      - Renders bullet lists up to 3 levels.
      - Replaces placeholder with a fitted image in the placeholder box.
    """

    # ZKB Corporate Style
    ZKB_BLUE = "003CD3"
    CORPORATE_FONT = "Frutiger for ZKB Light"

    # Adjust aliases if your template uses different localized names
    LAYOUT_NAME_ALIASES = {
        "title": ["Titelfolie", "Title Slide", "Titel", "Title"],
        "agenda": ["Agenda", "Inhaltsverzeichnis", "Table of Contents", "Title Only", "Blank"],
        "title_and_content": ["Title and Content", "Titel und Inhalt", "Title & Content", "Titel und Inhalt (1)"],
        "two_content": ["2 Spalten", "Two Content", "Zwei Inhalte", "Zwei Spalten"],
    }

    # Layout Mappings (Anchors in EMUs: 1 Inch = 914400 EMUs)
    # Coordinates from TemplateCreator.java:
    # 24, 148, 715, 374 (Title Slide Title)
    # 24, 23, 715, 20 (Title Slide Subtitle)
    # 24, 23, 802, 68 (Standard Slide Title)
    # 24, 148, 802, 374 (Standard Slide Content)
    # 24, 535, 100, 20 (Slide Number)
    
    # conversion factor: Java coordinates seem to be in points or a similar unit.
    # In POI, units for Rectangle2D.Double in setAnchor are usually points (1/72 inch).
    # python-pptx Pt(x) is 1/72 inch.

    def __init__(self, plan_path: str, output_path: str, files_dir: str, template_path: str):
        self.plan_path = plan_path
        self.output_path = output_path
        self.files_dir = files_dir
        self.template_path = template_path

        if not template_path or not os.path.exists(template_path):
            raise FileNotFoundError(f"template.pptx not found at: {template_path}")

        self.prs = Presentation(template_path)
        print(f"Using template: {template_path}")

        # Remove template slides marked with [layout: ...] in notes
        self._remove_layout_slides()

        # Resolve layout indices once, strictly
        self.layout_title = self._require_layout("title")
        self.layout_agenda = self._require_layout("agenda")
        self.layout_title_and_content = self._require_layout("title_and_content")
        self.layout_two_content = self._require_layout("two_content")

    def _remove_layout_slides(self) -> None:
        """
        Removes all slides from the presentation to start with a clean slate.
        """
        # We must be careful not to drop master/layout relationships.
        # Only drop slides.
        
        slide_rids = []
        for sld_id in self.prs.slides._sldIdLst.sldId_lst:
            slide_rids.append(sld_id.rId)
            
        # 1. Remove from slides collection (XML)
        for rId in slide_rids:
             # Find the sldId element with this rId
             for sld_id in list(self.prs.slides._sldIdLst.sldId_lst):
                 if sld_id.rId == rId:
                     self.prs.slides._sldIdLst.remove(sld_id)
                     break
            
        # 2. Drop slide relationships
        for rId in slide_rids:
            try:
                self.prs.part.drop_rel(rId)
            except Exception:
                pass
        
        print(f"Removed {len(slide_rids)} existing slides from presentation.")

    def _require_layout(self, logical_layout: str) -> int:
        names = self.LAYOUT_NAME_ALIASES.get(logical_layout, [])
        idx = _get_layout_index_by_any_name(self.prs, names)
        if idx is None:
            available = [layout.name for layout in self.prs.slide_layouts]
            raise RuntimeError(
                f"Required layout '{logical_layout}' not found in template.\n"
                f"Tried names: {names}\n"
                f"Available layouts: {available}"
            )
        return idx

    # ----------------------------
    # Parsing: YAML per slide
    # ----------------------------

    def parse_plan(self) -> List[SlideData]:
        """
        Slides are separated by '---'. Each block may begin with '# <title>'.
        Body is YAML.
        """
        with open(self.plan_path, "r", encoding="utf-8") as f:
            full_content = f.read()

        # Split by --- that is at the start of a line
        # Using re.MULTILINE to match ^---$
        blocks = re.split(r'^---$', full_content, flags=re.MULTILINE)
        
        # If no split occurred, try with just ---
        if len(blocks) == 1 and '---' in full_content:
             blocks = full_content.split('---')

        # If it starts with # at the very beginning of the file, we might not have a title header
        # But our slides start with #, so we should handle that.
        
        slides: List[SlideData] = []
        
        # Process title header if present
        first_block = blocks[0].strip()
        if first_block.startswith("%"):
            lines = first_block.splitlines()
            # Find the last line that starts with % to know where title header ends
            title_lines = []
            for line in lines:
                if line.startswith("%"):
                    title_lines.append(line[1:].strip())
                else:
                    break
            
            title = title_lines[0] if len(title_lines) > 0 else ""
            subtitle = title_lines[1] if len(title_lines) > 1 else ""
            author_date = title_lines[2] if len(title_lines) > 2 else ""
            if author_date:
                subtitle = f"{subtitle}\n{author_date}"
            
            slides.append({
                "layout": "title",
                "title": title,
                "subtitle": subtitle
            })
            
            # The rest of the first block after title lines should be treated as a new block if it contains content
            remaining_of_first = "\n".join(lines[len(title_lines):]).strip()
            if remaining_of_first:
                # If it starts with #, treat it as a slide. 
                # Pandoc often has the first slide immediately after title header without ---
                blocks[0] = remaining_of_first
            else:
                blocks = blocks[1:]
        
        for raw in blocks:
            # Further split blocks if they contain # at start of line (Pandoc style)
            # but only if they are not already separated by ---
            subblocks = re.split(r'^(?=# )', raw, flags=re.MULTILINE)
            for subblock in subblocks:
                block = subblock.strip()
                if not block:
                    continue

                lines = block.splitlines()
                title = ""
                yaml_body = block

                # If it starts with #, it's a title
                if lines and lines[0].strip().startswith("# "):
                    title = lines[0].strip()[2:].strip()
                    yaml_body = "\n".join(lines[1:]).strip()
                
                # Filter out Pandoc/Markdown specific markers that break YAML
                yaml_body = re.sub(r":::\s*notes.*?:::", "", yaml_body, flags=re.DOTALL)
                
                # Check for columns
                if "::: columns" in yaml_body:
                    data = {"layout": "two_content", "title": title}
                    cols = re.findall(r"::: column(.*?):::", yaml_body, flags=re.DOTALL)
                    if len(cols) >= 2:
                        left_col = cols[0].strip()
                        right_col = cols[1].strip()
                        
                        # Extract images from columns
                        left_img = re.search(r"!\[(.*?)\]\((.*?)\)", left_col)
                        right_img = re.search(r"!\[(.*?)\]\((.*?)\)", right_col)
                        
                        if left_img:
                            data["left"] = {"image": left_img.group(2), "caption": left_img.group(1)}
                        else:
                            data["left"] = left_col
                            
                        if right_img:
                            data["right"] = {"image": right_img.group(2), "caption": right_img.group(1)}
                        else:
                            data["right"] = right_col
                            
                        data["mode"] = "text_both" # will be upgraded in generate() if needed
                        slides.append(data)
                        continue

                yaml_body = re.sub(r":::\s*columns\s*", "", yaml_body)
                yaml_body = re.sub(r":::\s*column\s*", "", yaml_body)
                yaml_body = re.sub(r":::\s*", "", yaml_body)
                
                # If yaml_body is empty after title, it might be just a title slide or empty content
                if not yaml_body.strip():
                    slides.append({
                        "layout": "title_and_content",
                        "title": title,
                        "content": ""
                    })
                    continue

                # Best-effort cleanup to avoid common YAML mistakes
                yaml_body = _dedent_accidental_keys_in_block_scalars(yaml_body)
                yaml_body = _auto_indent_block_scalars(yaml_body)

                try:
                    data: SlideData = yaml.safe_load(yaml_body) or {}
                    # If we got a string instead of a dict, YAML thought it was a plain scalar
                    if not isinstance(data, dict):
                         data = {"content": yaml_body.strip()}
                except Exception as e:
                    # If YAML fails, check if it's just plain text or bullets
                    data = {"content": yaml_body.strip()}

                if title and "title" not in data:
                    data["title"] = title

                # If it's a simple markdown image link: ![alt](path)
                if "content" in data and not isinstance(data["content"], dict):
                    img_match = re.search(r"!\[(.*?)\]\((.*?)\)", str(data["content"]))
                    if img_match:
                        data["image"] = img_match.group(2)
                        data["caption"] = img_match.group(1)
                        data["content"] = re.sub(r"!\[.*?\]\(.*?\)", "", str(data["content"])).strip()

                # Normalize legacy layout aliases (layout: image_right => layout: two_content + mode)
                layout_raw = str(data.get("layout", "title_and_content") or "title_and_content").strip().lower()
                if layout_raw in ("image_right", "image_left", "text_both"):
                    data["layout"] = "two_content"
                    data["mode"] = layout_raw

                data.setdefault("layout", "title_and_content")
                data.setdefault("mode", None)

                for k in ("subtitle", "content", "left", "right", "image", "caption"):
                    if k in data and data[k] is None:
                        data[k] = ""

                slides.append(data)

        return slides

    # ----------------------------
    # Text (3 levels bullets)
    # ----------------------------

    def add_text_to_frame(self, text_frame, text: Union[str, List[str]]) -> None:
        """Render bullet list text with up to 3 levels; strips leading '-' markers."""
        text_frame.clear()
        text_frame.word_wrap = True

        lines = text if isinstance(text, list) else str(text or "").splitlines()
        non_empty = [l for l in lines if str(l).strip()]
        if not non_empty:
            return

        base_indent = min(len(l) - len(l.lstrip()) for l in non_empty) if non_empty else 0

        first = True
        for line in lines:
            if not str(line).strip():
                continue

            raw = str(line).rstrip("\n")
            stripped = raw.lstrip()
            indent = len(raw) - len(stripped)

            stripped = re.sub(r"^-\s+", "", stripped)

            level = max(0, min(((indent - base_indent) // 2), 2))

            p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            first = False

            p.level = level
            self._apply_formatted_text(p, stripped, level)

    def _apply_formatted_text(self, paragraph, text: str, level: int = 0) -> None:
        """Apply text to paragraph with support for bold formatting (**bold**)."""
        paragraph.font.name = self.CORPORATE_FONT
        
        # Determine font size based on level
        if level == 0:
            font_size = Pt(22)
        elif level == 1:
            font_size = Pt(18)
        elif level == -1:
            font_size = Pt(9)
        else:
            font_size = Pt(16)
            
        paragraph.font.size = font_size

        # Support bold text: **text**
        # Using a regex that captures the bold markers to split
        parts = re.split(r'(\*\*.*?\*\*)', text)
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                run = paragraph.add_run()
                run.text = part[2:-2]
                run.font.bold = True
            elif part:
                run = paragraph.add_run()
                run.text = part
            
            # Re-apply font properties to runs (python-pptx requires this for consistency)
            for run in paragraph.runs:
                run.font.name = self.CORPORATE_FONT
                run.font.size = font_size

    # ----------------------------
    # Images in placeholder box
    # ----------------------------

    def add_picture_in_placeholder_box(self, slide, placeholder, img_path: str) -> None:
        img_full = _abs_from_plan(self.plan_path, img_path)
        if not os.path.exists(img_full):
            # Try searching in files/images if not found directly
            alt_path = os.path.join(self.files_dir, "images", os.path.basename(img_path))
            if os.path.exists(alt_path):
                img_full = alt_path
            else:
                # Try relative to the script location too
                script_dir = os.path.dirname(os.path.abspath(__file__))
                alt_path2 = os.path.join(script_dir, img_path)
                if os.path.exists(alt_path2):
                    img_full = alt_path2
                else:
                    print(f"Warning: Image not found: {img_full}")
                    return

        left, top, box_w, box_h = placeholder.left, placeholder.top, placeholder.width, placeholder.height

        # Clear placeholder if it has text
        if hasattr(placeholder, 'has_text_frame') and placeholder.has_text_frame:
             placeholder.text_frame.clear()

        w, h = _fit_image_into_box(img_full, box_w, box_h)
        pic_left = left + (box_w - w) // 2
        pic_top = top + (box_h - h) // 2

        slide.shapes.add_picture(img_full, pic_left, pic_top, width=w, height=h)

    def add_table_to_placeholder(self, slide, placeholder, table_data: List[List[Any]]) -> None:
        """Render a table fitting inside the placeholder box."""
        if not table_data:
            return

        rows = len(table_data)
        cols = len(table_data[0]) if rows > 0 else 0
        if rows == 0 or cols == 0:
            return

        left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        
        # Instead of removing the placeholder, we just use its coordinates.
        # Removing can cause issues in some PPTX viewers or if we remove a required ph.
        # But we must ensure the placeholder is hidden or empty.
        if placeholder.has_text_frame:
             placeholder.text_frame.clear()

        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table

        # set column widths equally
        col_width = int(width / cols)
        for c in range(cols):
            table.columns[c].width = col_width

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = self.CORPORATE_FONT
                    paragraph.font.size = Pt(14)
                    if r == 0:  # Header
                        paragraph.font.bold = True
                        paragraph.alignment = PP_ALIGN.CENTER

    def add_footer_and_slide_number(self, slide, caption: str = "") -> None:
        """Add corporate footer and slide number with exact ZKB styling."""
        # Find slide number placeholder if it exists, otherwise create a textbox
        sn_ph = _find_placeholder(slide, placeholder_type=PP_PLACEHOLDER_TYPE.SLIDE_NUMBER)
        if sn_ph:
            # Enforce font and color in existing placeholder
            if hasattr(sn_ph, 'has_text_frame') and sn_ph.has_text_frame:
                tf = sn_ph.text_frame
                # We can't easily make it dynamic in python-pptx, but we can set the text
                slide_idx = list(self.prs.slides).index(slide) + 1
                tf.text = str(slide_idx)
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    p.font.name = self.CORPORATE_FONT
                    p.font.size = Pt(9)
                    p.font.color.rgb = RGBColor.from_string(self.ZKB_BLUE)
        else:
            # 24, 535, 100, 20 (in points)
            sn_box = slide.shapes.add_textbox(Pt(24), Pt(535), Pt(100), Pt(20))
            tf = sn_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            
            slide_idx = list(self.prs.slides).index(slide) + 1
            p.text = str(slide_idx)
            p.font.name = self.CORPORATE_FONT
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor.from_string(self.ZKB_BLUE)

        # Footer / Caption
        footer_text = caption or "AngularAI - Softwareentwicklung mit AI"
        
        # Look for footer placeholder
        f_ph = _find_placeholder(slide, placeholder_type=PP_PLACEHOLDER_TYPE.FOOTER)
        if f_ph:
            if hasattr(f_ph, 'has_text_frame') and f_ph.has_text_frame:
                ftf = f_ph.text_frame
                ftf.clear()
                self._apply_formatted_text(ftf.paragraphs[0], footer_text, level=-1) # level -1 for footer size
                for p in ftf.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    p.font.color.rgb = RGBColor.from_string("808080")
        else:
            # Position it next to slide number or centered
            # template-company: Footer is at (48.16, 535.41, 777.62, 15.01) pts
            f_box = slide.shapes.add_textbox(Pt(60), Pt(535), Pt(600), Pt(20))
            ftf = f_box.text_frame
            ftf.clear()
            self._apply_formatted_text(ftf.paragraphs[0], footer_text, level=-1)
            for p in ftf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.font.color.rgb = RGBColor.from_string("808080") # Grey for footer

    # ----------------------------
    # Slide builders (template layouts only)
    # ----------------------------

    def add_title_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[self.layout_title]
        slide = self.prs.slides.add_slide(layout)

        # In template.pptx, we already have the background image if added by TemplateCreator.
        # Ensure we don't add it again if it's already in the layout.
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("title", "")
            # Force position: (24, 148, 715, 374) pts from template-company
            title_shape.left = Pt(24)
            title_shape.top = Pt(148)
            title_shape.width = Pt(715)
            title_shape.height = Pt(374)
            
            p = title_shape.text_frame.paragraphs[0]
            p.font.name = self.CORPORATE_FONT
            p.font.size = Pt(36)
            p.font.color.rgb = RGBColor.from_string(self.ZKB_BLUE)
            p.font.bold = True
        
        subtitle = str(slide_data.get("subtitle", "") or "").strip()
        if subtitle:
            # Anchor: (24.08, 22.66, 715, 67.9) pts from template-company
            ph = _find_placeholder(slide, placeholder_type=PP_PLACEHOLDER_TYPE.SUBTITLE)
            if not ph:
                 # fallback to searching for a shape with "Untertitel" or similar if PH not found
                 pass
                 
            if ph and ph.has_text_frame:
                ph.left = Pt(24)
                ph.top = Pt(23)
                ph.width = Pt(715)
                ph.height = Pt(68)
                ph.text = subtitle
                p = ph.text_frame.paragraphs[0]
                p.font.name = self.CORPORATE_FONT
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor.from_string(self.ZKB_BLUE)
        
        # Add image to title slide ONLY if it's NOT the background (already in template)
        img_path = str(slide_data.get("image", "") or "").strip()
        if img_path and "AiRace.png" not in img_path:
            ph = _find_placeholder(slide, placeholder_type=PP_PLACEHOLDER_TYPE.OBJECT)
            if ph:
                self.add_picture_in_placeholder_box(slide, ph, img_path)
            else:
                img_full = _abs_from_plan(self.plan_path, img_path)
                if os.path.exists(img_full):
                    w, h = _fit_image_into_box(img_full, Inches(4), Inches(3))
                    slide.shapes.add_picture(img_full, (self.prs.slide_width - w) // 2, Inches(4), width=w, height=h)

        self.add_footer_and_slide_number(slide)

    def add_agenda_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[self.layout_agenda]
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "Agenda")
            for p in slide.shapes.title.text_frame.paragraphs:
                p.font.name = self.CORPORATE_FONT
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor.from_string(self.ZKB_BLUE)
                p.font.bold = True

        body = None
        for ph in slide.placeholders:
            if slide.shapes.title and ph == slide.shapes.title:
                continue
            if ph.has_text_frame:
                body = ph
                break
        if body:
            self.add_text_to_frame(body.text_frame, slide_data.get("content", ""))
        
        self.add_footer_and_slide_number(slide)

    def add_title_and_content_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[self.layout_title_and_content]
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")
            # Force position: (24, 23, 802, 68) in points
            slide.shapes.title.left = Pt(24)
            slide.shapes.title.top = Pt(23)
            slide.shapes.title.width = Pt(802)
            slide.shapes.title.height = Pt(68)
            for p in slide.shapes.title.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT # Ensure it's in the corner
                p.font.name = self.CORPORATE_FONT
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor.from_string(self.ZKB_BLUE)
                p.font.bold = True

        # Check for image and content
        content = str(slide_data.get("content", "") or "").strip()
        img_path = str(slide_data.get("image", "") or "").strip()
        
        # If we have both, we should have been upgraded to two_content already.
        # But if we just have an image, place it in the content area.
        
        body = None
        for ph in slide.placeholders:
            if slide.shapes.title and ph == slide.shapes.title:
                continue
            if ph.placeholder_format.type in (PP_PLACEHOLDER_TYPE.BODY, PP_PLACEHOLDER_TYPE.OBJECT):
                body = ph
                break
        
        if body:
            # Force body position if needed: (24, 148, 802, 374)
            body.left = Pt(24)
            body.top = Pt(148)
            body.width = Pt(802)
            body.height = Pt(374)
            
            if img_path:
                self.add_picture_in_placeholder_box(slide, body, img_path)
            elif content:
                self.add_text_to_frame(body.text_frame, content)
        elif img_path:
            # Fallback if no body placeholder
            img_full = _abs_from_plan(self.plan_path, img_path)
            if os.path.exists(img_full):
                w, h = _fit_image_into_box(img_full, Pt(802), Pt(374))
                slide.shapes.add_picture(img_full, Pt(24), Pt(148), width=w, height=h)
        
        self.add_footer_and_slide_number(slide)

    def _parse_side(self, side: Any) -> Dict[str, Any]:
        """Normalize a 'left'/'right' side into either {'text':...} or {'image':..., 'caption':...}."""
        if isinstance(side, dict):
            out = dict(side)
            if "img" in out and "image" not in out:
                out["image"] = out.pop("img")
            # If mapping has no explicit text key but does have 'text' absent, keep as-is
            return out
        if isinstance(side, str):
            s = side.strip()
            if not s:
                return {"text": ""}
            return {"text": side}
        return {"text": str(side)}

    def add_two_content_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[self.layout_two_content]
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")
            # Force position: (24, 23, 802, 68) in points
            slide.shapes.title.left = Pt(24)
            slide.shapes.title.top = Pt(23)
            slide.shapes.title.width = Pt(802)
            slide.shapes.title.height = Pt(68)
            slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            slide.shapes.title.text_frame.paragraphs[0].font.name = self.CORPORATE_FONT
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.ZKB_BLUE)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        # Robustly find the two content placeholders in the layout.
        candidates = []
        for ph in slide.placeholders:
            try:
                pht = ph.placeholder_format.type
            except Exception:
                continue

            if pht in (PP_PLACEHOLDER_TYPE.OBJECT, PP_PLACEHOLDER_TYPE.BODY):
                if slide.shapes.title is not None and ph == slide.shapes.title:
                    continue
                candidates.append(ph)

        if len(candidates) < 2:
            for ph in slide.placeholders:
                if slide.shapes.title is not None and ph == slide.shapes.title:
                    continue
                if getattr(ph, "has_text_frame", False):
                    candidates.append(ph)
                if len(candidates) >= 2:
                    break

        if len(candidates) < 2:
            raise RuntimeError("Two-content layout must provide two content placeholders (OBJECT/BODY).")

        candidates.sort(key=lambda p: p.left)
        left_ph, right_ph = candidates[0], candidates[1]
        
        # Force positions based on TemplateCreator.java
        # Content 1: (24, 148, 380, 374) - slightly adjusted to 470 in comments but 380 in code
        # Actually TemplateCreator says:
        # Content 1: [Shape: "Inhaltsplatzhalter 4", ID: 5, PH Type: obj, Anchor: (24.08, 148.44, 470.7, 373.9)]
        # Content 2: [Shape: "Inhaltsplatzhalter 5", ID: 6, PH Type: obj, Anchor: (517.45, 148.44, 470.7, 373.9)]
        left_ph.left = Pt(24)
        left_ph.top = Pt(148)
        left_ph.width = Pt(470)
        left_ph.height = Pt(374)
        
        right_ph.left = Pt(517)
        right_ph.top = Pt(148)
        right_ph.width = Pt(470)
        right_ph.height = Pt(374)

        mode = (slide_data.get("mode") or "").strip().lower() or "text_both"

        left_side = self._parse_side(slide_data.get("left", ""))
        right_side = self._parse_side(slide_data.get("right", ""))
        
        caption = ""

        def coerce_image(side: Dict[str, Any]) -> Dict[str, Any]:
            if "image" in side:
                return side
            txt = str(side.get("text", "") or "").strip()
            if txt.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
                return {"image": txt}
            return side

        if mode == "image_left":
            left_side = coerce_image(left_side)
        elif mode == "image_right":
            right_side = coerce_image(right_side)

                        # Render left
        if "image" in left_side and str(left_side["image"]).strip():
            self.add_picture_in_placeholder_box(slide, left_ph, str(left_side["image"]).strip())
            # Check for caption in the side mapping
            side_caption = str(left_side.get("caption", "") or "").strip()
            if side_caption:
                caption = side_caption
        elif "table" in left_side:
            self.add_table_to_placeholder(slide, left_ph, left_side["table"])
        else:
            self.add_text_to_frame(left_ph.text_frame, left_side.get("text", ""))

                        # Render right
        if "image" in right_side and str(right_side["image"]).strip():
            self.add_picture_in_placeholder_box(slide, right_ph, str(right_side["image"]).strip())
            # Check for caption in the side mapping
            side_caption = str(right_side.get("caption", "") or "").strip()
            if side_caption:
                caption = side_caption
        elif "table" in right_side:
            self.add_table_to_placeholder(slide, right_ph, right_side["table"])
        else:
            self.add_text_to_frame(right_ph.text_frame, right_side.get("text", ""))

        if caption:
            self.add_footer_and_slide_number(slide, caption=caption)
        else:
            self.add_footer_and_slide_number(slide)

    # ----------------------------
    # Generation: enforce template-only layouts + upgrade rules
    # ----------------------------

    def generate(self) -> None:
        # TemplateCreator already produces a clean template.pptx without slides.
        # However, it seems some phantom slides still exist or are added.
        # We start with a fresh new presentation if the template has slides
        # to ensure NO phantom slides.
        if len(self.prs.slides) > 0:
            print(f"Warning: Template contains {len(self.prs.slides)} slides. Removing them...")
            slide_rids = []
            for sld_id in self.prs.slides._sldIdLst.sldId_lst:
                slide_rids.append(sld_id.rId)
                
            for rId in slide_rids:
                for sld_id in list(self.prs.slides._sldIdLst.sldId_lst):
                    if sld_id.rId == rId:
                        self.prs.slides._sldIdLst.remove(sld_id)
                        break
                try:
                    self.prs.part.drop_rel(rId)
                except Exception:
                    pass

        slides_data = self.parse_plan()

        for data in slides_data:
            layout_type = str(data.get("layout", "title_and_content") or "title_and_content").strip().lower()

            # accept legacy aliases at generation time too
            if layout_type in ("image_right", "image_left", "text_both"):
                data = dict(data)
                data["layout"] = "two_content"
                data["mode"] = layout_type
                layout_type = "two_content"

            # Upgrade: title_and_content with both content+image => two_content (image_right by default)
            if layout_type == "title_and_content":
                has_text = bool(str(data.get("content", "") or "").strip())
                has_img = bool(str(data.get("image", "") or "").strip())
                has_table = bool(data.get("table"))
                if has_text and (has_img or has_table):
                    layout_type = "two_content"
                    data = dict(data)
                    data["mode"] = data.get("mode") or ("image_right" if has_img else "text_both")
                    data["left"] = {"text": data.get("content", "")}
                    if has_img:
                        data["right"] = {"image": data.get("image", ""), "caption": data.get("caption", "")}
                    else:
                        data["right"] = {"table": data.get("table")}

            print(f"Adding slide: {data.get('title', 'Untitled')} (Layout: {layout_type})")
            if layout_type == "title":
                self.add_title_slide(data)
            elif layout_type == "agenda":
                self.add_agenda_slide(data)
            elif layout_type == "two_content":
                self.add_two_content_slide(data)
            elif layout_type == "title_and_content":
                self.add_title_and_content_slide(data)
            else:
                raise RuntimeError(
                    f"Unknown layout '{layout_type}'. Allowed: title, agenda, title_and_content, two_content "
                    f"(or legacy: image_left/image_right/text_both)."
                )

        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        try:
            self.prs.save(self.output_path)
            print(f"Presentation saved to: {self.output_path}")
        except PermissionError:
            alt_path = self.output_path.replace(".pptx", "_alt.pptx")
            print(f"Warning: Could not save to {self.output_path} (file might be open). Saving to {alt_path} instead.")
            self.prs.save(alt_path)


if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.abspath(__file__))

    plan = os.path.join(base_dir, "presentation-slides.md")
    output = os.path.join(base_dir, "SoftwareEntwicklungAI.pptx")
    files = os.path.join(base_dir, "files")
    template = os.path.join(base_dir, "template.pptx")

    gen = PresentationGenerator(plan, output, files, template)
    gen.generate()
