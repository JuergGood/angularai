import os
import re
from typing import Any, Dict, List, Optional, Tuple, Union

import yaml
from PIL import Image

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, PP_PLACEHOLDER_TYPE
from pptx.util import Inches, Pt


SlideData = Dict[str, Any]


# ----------------------------
# Path + media helpers
# ----------------------------

def _abs_from_plan(plan_path: str, maybe_rel: str) -> str:
    """Resolve path relative to the markdown plan file."""
    if not maybe_rel:
        return maybe_rel
    if os.path.isabs(maybe_rel):
        return maybe_rel
    return os.path.normpath(os.path.join(os.path.dirname(plan_path), maybe_rel))


def _fit_image_into_box(img_path: str, box_w_emu: int, box_h_emu: int) -> Tuple[int, int]:
    """Return (w_emu, h_emu) preserving aspect ratio to fit within the box (contain)."""
    with Image.open(img_path) as im:
        w_px, h_px = im.size
    if w_px <= 0 or h_px <= 0:
        return box_w_emu, box_h_emu

    img_ar = w_px / h_px
    box_ar = box_w_emu / box_h_emu

    if img_ar >= box_ar:
        w = box_w_emu
        h = int(box_w_emu / img_ar)
    else:
        h = box_h_emu
        w = int(box_h_emu * img_ar)
    return w, h


# ----------------------------
# YAML robustness helpers
# ----------------------------

_TOP_LEVEL_KEY_RE = re.compile(r"^[A-Za-z_]\w*:\s*(.*)$")
_BLOCK_SCALAR_START_RE = re.compile(r"^[A-Za-z_]\w*:\s*\|\s*$")


def _auto_indent_block_scalars(yaml_body: str) -> str:
    """
    Makes YAML more forgiving for 'key: |' blocks by ensuring subsequent lines
    are indented at least 2 spaces until the next top-level key.

    Fixes common mistakes like:
      content: |
      - item
    to:
      content: |
        - item

    Note: This does not fix structurally incorrect indentation where a new key
    (e.g. 'right:') was accidentally written inside a block scalar. For that,
    see _dedent_accidental_keys_in_block_scalars().
    """
    lines = yaml_body.splitlines()
    out: List[str] = []
    in_block = False

    for line in lines:
        if _BLOCK_SCALAR_START_RE.match(line):
            in_block = True
            out.append(line)
            continue

        # end block when we see a new top-level key (no leading spaces)
        if in_block and line and not line.startswith((" ", "\t")) and _TOP_LEVEL_KEY_RE.match(line):
            in_block = False

        if in_block:
            if line.strip() and not line.startswith((" ", "\t")):
                out.append("  " + line)
            else:
                out.append(line)
        else:
            out.append(line)

    return "\n".join(out)


def _dedent_accidental_keys_in_block_scalars(yaml_body: str, keys: List[str] = ["right:", "left:", "image:", "caption:", "text:", "mode:", "layout:"]) -> str:
    """
    Best-effort fix for mistakes like:

      left: |
        - bullet
        right:
          image: x.png

    where 'right:' was intended as a top-level key but is indented within the left block scalar.

    Strategy:
      - Track when we're inside a block scalar.
      - If we encounter a line with two+ leading spaces that (after strip) begins with one of the known keys,
        we dedent it to column 0 and also dedent subsequent more-indented lines (the nested mapping) by 2 spaces
        until we hit a non-indented top-level key.

    This is conservative and only targets known keys.
    """
    key_prefixes = tuple(keys)
    lines = yaml_body.splitlines()
    out: List[str] = []
    in_block = False
    pending_dedent = False

    for i, line in enumerate(lines):
        if _BLOCK_SCALAR_START_RE.match(line):
            in_block = True
            pending_dedent = False
            out.append(line)
            continue

        # If we hit a new top-level key, we are no longer in the block scalar
        if in_block and line and not line.startswith((" ", "\t")) and _TOP_LEVEL_KEY_RE.match(line):
            in_block = False
            pending_dedent = False
            out.append(line)
            continue

        if in_block:
            stripped = line.lstrip()
            indent = len(line) - len(stripped)

            if indent >= 2 and stripped.startswith(key_prefixes):
                # Dedent this line to become a top-level key
                out.append(stripped)
                pending_dedent = True
                continue

            if pending_dedent:
                # Dedent nested lines by 2 spaces if possible
                if line.startswith("  "):
                    out.append(line[2:])
                    continue
                else:
                    # If we hit a non-indented line inside pending_dedent, stop dedenting
                    pending_dedent = False
                    out.append(line)
                    continue

            out.append(line)
        else:
            out.append(line)

    return "\n".join(out)


# ----------------------------
# PPTX layout + placeholder helpers
# ----------------------------

def _norm_layout_name(name: str) -> str:
    return re.sub(r"\s+", " ", (name or "").strip().lower())


def _get_layout_index_by_name(prs: Presentation, wanted: str) -> Optional[int]:
    wn = _norm_layout_name(wanted)
    for i, layout in enumerate(prs.slide_layouts):
        if _norm_layout_name(layout.name) == wn:
            return i
    return None


def _get_layout_index_by_any_name(prs: Presentation, wanted_names: List[str]) -> Optional[int]:
    for n in wanted_names:
        idx = _get_layout_index_by_name(prs, n)
        if idx is not None:
            return idx
    return None


def _find_placeholder(slide, idx: Optional[int] = None, placeholder_type: Optional[int] = None):
    """Best-effort placeholder lookup by index or placeholder type."""
    if idx is not None:
        try:
            return slide.placeholders[idx]
        except Exception:
            return None
    if placeholder_type is not None:
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.type == placeholder_type:
                    return ph
            except Exception:
                continue
    return None


def _remove_shape(shape) -> None:
    try:
        el = shape.element
        el.getparent().remove(el)
    except Exception:
        pass


# ----------------------------
# Generator
# ----------------------------

class PresentationGenerator:
    """
    Key properties:
      - Uses layouts ONLY from template.pptx by layout name (no index fallbacks).
      - Parses each slide block as YAML.
      - Adds layout: agenda.
      - Unifies all two-content variants into:
            layout: two_content
            mode: image_left | image_right | text_both
        Additionally supports legacy aliases:
            layout: image_left / image_right / text_both
      - Renders bullet lists up to 3 levels.
      - Replaces placeholder with a fitted image in the placeholder box.
    """

    # Adjust aliases if your template uses different localized names
    LAYOUT_NAME_ALIASES = {
        "title": ["Titelfolie", "Title Slide", "Titel", "Title"],
        "agenda": ["Agenda", "Inhaltsverzeichnis", "Table of Contents"],
        "title_and_content": ["Title and Content", "Titel und Inhalt", "Title & Content", "Titel und Inhalt (1)"],
        "two_content": ["2 Spalten", "Two Content", "Zwei Inhalte", "Zwei Spalten"],
    }

    def __init__(self, plan_path: str, output_path: str, files_dir: str, template_path: str):
        self.plan_path = plan_path
        self.output_path = output_path
        self.files_dir = files_dir
        self.template_path = template_path

        if not template_path or not os.path.exists(template_path):
            raise FileNotFoundError(f"template.pptx not found at: {template_path}")

        self.prs = Presentation(template_path)
        print(f"Using template: {template_path}")

        # Remove template slides marked with [layout: ...] in notes
        self._remove_layout_slides()

        # Resolve layout indices once, strictly
        self.layout_title = self._require_layout("title")
        self.layout_agenda = self._require_layout("agenda")
        self.layout_title_and_content = self._require_layout("title_and_content")
        self.layout_two_content = self._require_layout("two_content")

    def _remove_layout_slides(self) -> None:
        """
        Removes slides from the presentation that are tagged as 'layout' templates
        in their slide notes (e.g., [layout: title]).
        """
        slides_to_remove = []
        for slide in self.prs.slides:
            try:
                notes = slide.notes_slide.notes_text_frame.text
                if "[layout:" in notes.lower():
                    slides_to_remove.append(slide)
            except Exception:
                continue

        for slide in slides_to_remove:
            slide_id = slide.slide_id
            
            # 1. Remove from slides collection (XML)
            xml_slides = self.prs.slides._sldIdLst
            for sld_id in xml_slides.sldId_lst:
                if sld_id.id == slide_id:
                    xml_slides.remove(sld_id)
                    break
            
            # 2. Remove the slide part from the package to avoid "Duplicate name" warnings on save
            slide_part = slide.part
            rel_id = None
            for rId, rel in self.prs.part.rels.items():
                if rel.target_part == slide_part:
                    rel_id = rId
                    break
            if rel_id:
                self.prs.part.drop_rel(rel_id)

        if slides_to_remove:
            print(f"Removed {len(slides_to_remove)} layout template slides from presentation.")

    def _require_layout(self, logical_layout: str) -> int:
        names = self.LAYOUT_NAME_ALIASES.get(logical_layout, [])
        idx = _get_layout_index_by_any_name(self.prs, names)
        if idx is None:
            available = [layout.name for layout in self.prs.slide_layouts]
            raise RuntimeError(
                f"Required layout '{logical_layout}' not found in template.\n"
                f"Tried names: {names}\n"
                f"Available layouts: {available}"
            )
        return idx

    # ----------------------------
    # Parsing: YAML per slide
    # ----------------------------

    def parse_plan(self) -> List[SlideData]:
        """
        Slides are separated by '---'. Each block may begin with '# <title>'.
        Body is YAML.

        Supports:
          - content: |  <multi-line bullets>
          - left: | / right: | (multi-line)
          - left/right as mapping: {image:, caption:, text: }
          - layout aliases: image_left/image_right/text_both
        """
        with open(self.plan_path, "r", encoding="utf-8") as f:
            content = f.read()

        slides: List[SlideData] = []
        for raw in content.split("---"):
            block = raw.strip()
            if not block:
                continue

            lines = block.splitlines()
            title = ""
            yaml_body = block

            if lines and lines[0].startswith("# "):
                title = lines[0][2:].strip()
                yaml_body = "\n".join(lines[1:]).strip()

            # Best-effort cleanup to avoid common YAML mistakes
            yaml_body = _dedent_accidental_keys_in_block_scalars(yaml_body)
            yaml_body = _auto_indent_block_scalars(yaml_body)

            try:
                data: SlideData = yaml.safe_load(yaml_body) or {}
            except Exception as e:
                raise RuntimeError(f"YAML parse failed for slide '{title}': {e}") from e

            if title and "title" not in data:
                data["title"] = title

            # Normalize legacy layout aliases (layout: image_right => layout: two_content + mode)
            layout_raw = str(data.get("layout", "title_and_content") or "title_and_content").strip().lower()
            if layout_raw in ("image_right", "image_left", "text_both"):
                data["layout"] = "two_content"
                data["mode"] = layout_raw

            data.setdefault("layout", "title_and_content")
            data.setdefault("mode", None)

            for k in ("subtitle", "content", "left", "right", "image", "caption"):
                if k in data and data[k] is None:
                    data[k] = ""

            slides.append(data)

        return slides

    # ----------------------------
    # Text (3 levels bullets)
    # ----------------------------

    def add_text_to_frame(self, text_frame, text: Union[str, List[str]]) -> None:
        """Render bullet list text with up to 3 levels; strips leading '-' markers."""
        text_frame.clear()

        lines = text if isinstance(text, list) else str(text or "").splitlines()
        non_empty = [l for l in lines if str(l).strip()]
        if not non_empty:
            return

        base_indent = min(len(l) - len(l.lstrip()) for l in non_empty)

        first = True
        for line in lines:
            if not str(line).strip():
                continue

            raw = str(line).rstrip("\n")
            stripped = raw.lstrip()
            indent = len(raw) - len(stripped)

            stripped = re.sub(r"^-\s+", "", stripped)

            level = max(0, min(((indent - base_indent) // 2), 2))

            p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            first = False

            p.text = stripped
            p.level = level

            if level == 0:
                p.font.size = Pt(22)
            elif level == 1:
                p.font.size = Pt(18)
            else:
                p.font.size = Pt(16)

    # ----------------------------
    # Images in placeholder box
    # ----------------------------

    def add_picture_in_placeholder_box(self, slide, placeholder, img_path: str) -> None:
        img_full = _abs_from_plan(self.plan_path, img_path)
        if not os.path.exists(img_full):
            raise FileNotFoundError(f"Image not found: {img_full}")

        left, top, box_w, box_h = placeholder.left, placeholder.top, placeholder.width, placeholder.height

        _remove_shape(placeholder)

        w, h = _fit_image_into_box(img_full, box_w, box_h)
        pic_left = left + int((box_w - w) / 2)
        pic_top = top + int((box_h - h) / 2)

        slide.shapes.add_picture(img_full, pic_left, pic_top, width=w, height=h)

    # ----------------------------
    # Slide builders (template layouts only)
    # ----------------------------

    def add_title_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[self.layout_title]
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")
        else:
            for ph in slide.placeholders:
                if ph.has_text_frame:
                    ph.text = slide_data.get("title", "")
                    break

        subtitle = str(slide_data.get("subtitle", "") or "").strip()
        if subtitle:
            ph = _find_placeholder(slide, idx=1)
            if ph and ph.has_text_frame:
                ph.text = subtitle

    def add_agenda_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[self.layout_agenda]
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "Agenda")

        body = None
        for ph in slide.placeholders:
            if slide.shapes.title and ph == slide.shapes.title:
                continue
            if ph.has_text_frame:
                body = ph
                break
        if body:
            self.add_text_to_frame(body.text_frame, slide_data.get("content", ""))

    def add_title_and_content_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[self.layout_title_and_content]
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")

        body = None
        for ph in slide.placeholders:
            if slide.shapes.title and ph == slide.shapes.title:
                continue
            if ph.has_text_frame:
                body = ph
                break
        if body:
            self.add_text_to_frame(body.text_frame, slide_data.get("content", ""))

    def _parse_side(self, side: Any) -> Dict[str, Any]:
        """Normalize a 'left'/'right' side into either {'text':...} or {'image':..., 'caption':...}."""
        if isinstance(side, dict):
            out = dict(side)
            if "img" in out and "image" not in out:
                out["image"] = out.pop("img")
            # If mapping has no explicit text key but does have 'text' absent, keep as-is
            return out
        if isinstance(side, str):
            s = side.strip()
            if not s:
                return {"text": ""}
            return {"text": side}
        return {"text": str(side)}

    def add_two_content_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[self.layout_two_content]
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")

        # Robustly find the two content placeholders in the layout.
        # Some templates don't use placeholder idx 1 and 2, so we select the first two
        # OBJECT/BODY placeholders (excluding title/date/footer/slide number) and order them by x-position.
        candidates = []
        for ph in slide.placeholders:
            try:
                pht = ph.placeholder_format.type
            except Exception:
                continue

            if pht in (PP_PLACEHOLDER_TYPE.OBJECT, PP_PLACEHOLDER_TYPE.BODY):
                # Exclude title placeholder if it is also of a compatible type (rare)
                if slide.shapes.title is not None and ph == slide.shapes.title:
                    continue
                candidates.append(ph)

        if len(candidates) < 2:
            # Fall back to any non-title text placeholders
            for ph in slide.placeholders:
                if slide.shapes.title is not None and ph == slide.shapes.title:
                    continue
                if getattr(ph, "has_text_frame", False):
                    candidates.append(ph)
                if len(candidates) >= 2:
                    break

        if len(candidates) < 2:
            raise RuntimeError("Two-content layout must provide two content placeholders (OBJECT/BODY).")

        candidates.sort(key=lambda p: p.left)
        left_ph, right_ph = candidates[0], candidates[1]

        mode = (slide_data.get("mode") or "").strip().lower() or "text_both"

        left_side = self._parse_side(slide_data.get("left", ""))
        right_side = self._parse_side(slide_data.get("right", ""))

        def coerce_image(side: Dict[str, Any]) -> Dict[str, Any]:
            if "image" in side:
                return side
            txt = str(side.get("text", "") or "").strip()
            if txt.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
                return {"image": txt}
            return side

        if mode == "image_left":
            left_side = coerce_image(left_side)
        elif mode == "image_right":
            right_side = coerce_image(right_side)

        # Render left
        if "image" in left_side and str(left_side["image"]).strip():
            self.add_picture_in_placeholder_box(slide, left_ph, str(left_side["image"]).strip())
        else:
            self.add_text_to_frame(left_ph.text_frame, left_side.get("text", ""))

        # Render right
        if "image" in right_side and str(right_side["image"]).strip():
            self.add_picture_in_placeholder_box(slide, right_ph, str(right_side["image"]).strip())
        else:
            self.add_text_to_frame(right_ph.text_frame, right_side.get("text", ""))

        caption = str(right_side.get("caption", "") or left_side.get("caption", "") or "").strip()
        if caption:
            footer = _find_placeholder(slide, placeholder_type=PP_PLACEHOLDER.FOOTER)
            if footer and footer.has_text_frame:
                footer.text = caption
            else:
                box = slide.shapes.add_textbox(Inches(7.0), Inches(6.9), Inches(6.0), Inches(0.4))
                tf = box.text_frame
                tf.text = caption
                tf.paragraphs[0].font.size = Pt(12)

    # ----------------------------
    # Generation: enforce template-only layouts + upgrade rules
    # ----------------------------

    def generate(self) -> None:
        slides_data = self.parse_plan()

        for data in slides_data:
            layout_type = str(data.get("layout", "title_and_content") or "title_and_content").strip().lower()

            # accept legacy aliases at generation time too
            if layout_type in ("image_right", "image_left", "text_both"):
                data = dict(data)
                data["layout"] = "two_content"
                data["mode"] = layout_type
                layout_type = "two_content"

            # Upgrade: title_and_content with both content+image => two_content (image_right by default)
            if layout_type == "title_and_content":
                has_text = bool(str(data.get("content", "") or "").strip())
                has_img = bool(str(data.get("image", "") or "").strip())
                if has_text and has_img:
                    layout_type = "two_content"
                    data = dict(data)
                    data["mode"] = data.get("mode") or "image_right"
                    data["left"] = {"text": data.get("content", "")}
                    data["right"] = {"image": data.get("image", ""), "caption": data.get("caption", "")}

            if layout_type == "title":
                self.add_title_slide(data)
            elif layout_type == "agenda":
                self.add_agenda_slide(data)
            elif layout_type == "two_content":
                self.add_two_content_slide(data)
            elif layout_type == "title_and_content":
                self.add_title_and_content_slide(data)
            else:
                raise RuntimeError(
                    f"Unknown layout '{layout_type}'. Allowed: title, agenda, title_and_content, two_content "
                    f"(or legacy: image_left/image_right/text_both)."
                )

        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        try:
            self.prs.save(self.output_path)
            print(f"Presentation saved to: {self.output_path}")
        except PermissionError:
            alt_path = self.output_path.replace(".pptx", "_alt.pptx")
            print(f"Warning: Could not save to {self.output_path} (file might be open). Saving to {alt_path} instead.")
            self.prs.save(alt_path)


if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.abspath(__file__))

    plan = os.path.join(base_dir, "presentation-slides-de.md")
    output = os.path.join(base_dir, "generated", "SoftwareDevelopmentWithAI.pptx")
    files = os.path.join(base_dir, "files")
    template = os.path.join(base_dir, "template.pptx")

    gen = PresentationGenerator(plan, output, files, template)
    gen.generate()
