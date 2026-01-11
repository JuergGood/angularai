import os
import re
from typing import Any, Dict, List, Optional, Union

import yaml
from PIL import Image

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt


SlideData = Dict[str, Any]


def _abs_from_plan(plan_path: str, maybe_rel: str) -> str:
    """Resolve path relative to the markdown plan file."""
    if not maybe_rel:
        return maybe_rel
    if os.path.isabs(maybe_rel):
        return maybe_rel
    return os.path.normpath(os.path.join(os.path.dirname(plan_path), maybe_rel))


def _fit_image_into_box(img_path: str, box_w_emu: int, box_h_emu: int) -> (int, int):
    """Return (w_emu, h_emu) preserving aspect ratio to fit within the box (contain)."""
    with Image.open(img_path) as im:
        w_px, h_px = im.size
    if w_px <= 0 or h_px <= 0:
        return box_w_emu, box_h_emu

    img_ar = w_px / h_px
    box_ar = box_w_emu / box_h_emu

    if img_ar >= box_ar:
        w = box_w_emu
        h = int(box_w_emu / img_ar)
    else:
        h = box_h_emu
        w = int(box_h_emu * img_ar)
    return w, h


def _find_placeholder(slide, idx: Optional[int] = None, placeholder_type: Optional[int] = None):
    """Best-effort placeholder lookup by index or placeholder type."""
    if idx is not None:
        try:
            return slide.placeholders[idx]
        except Exception:
            return None
    if placeholder_type is not None:
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.type == placeholder_type:
                    return ph
            except Exception:
                continue
    return None


class PresentationGenerator:
    """
    Proposal A: Parse each slide block as YAML (robust; supports nested right:{image,caption})
    Proposal B: Place images using placeholder bounding box, keep aspect ratio, center
    Proposal C: Auto-layout upgrade: if title_and_content has both content+image -> render as Two Content
    """

    def __init__(self, plan_path: str, output_path: str, files_dir: str, template_path: Optional[str] = None):
        self.plan_path = plan_path
        self.output_path = output_path
        self.files_dir = files_dir

        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            print(f"Using template: {template_path}")
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.33)
            self.prs.slide_height = Inches(7.5)
            if template_path:
                print(f"Warning: Template not found at {template_path}. Using default.")

    # ----------------------------
    # Proposal A: YAML parsing
    # ----------------------------

    def parse_plan(self) -> List[SlideData]:
        """
        Your plan format is: blocks separated by '---'
          # Title
          layout: two_content
          left: |
            - Bullet
              - Sub bullet
          right:
            image: files/generated/foo.png
            caption: Optional
            text: Optional bullets
        """
        with open(self.plan_path, "r", encoding="utf-8") as f:
            content = f.read()

        slides: List[SlideData] = []
        for raw in content.split("---"):
            block = raw.strip()
            if not block:
                continue

            lines = block.splitlines()
            title = ""
            yaml_body = block

            if lines and lines[0].startswith("# "):
                title = lines[0][2:].strip()
                yaml_body = "\n".join(lines[1:]).strip()

            data: SlideData = {}
            if yaml_body:
                try:
                    data = yaml.safe_load(yaml_body) or {}
                except Exception as e:
                    # Keep your legacy parsing as a fallback so generation never breaks
                    print(f"Warning: YAML parse failed for slide '{title}': {e}. Falling back to legacy parsing.")
                    data = self._legacy_parse_block(yaml_body)

            if title and "title" not in data:
                data["title"] = title

            data.setdefault("layout", "title_and_content")
            data.setdefault("subtitle", "")
            data.setdefault("content", "")
            data.setdefault("left", "")
            data.setdefault("right", "")
            data.setdefault("image", "")

            slides.append(data)

        return slides

    def _legacy_parse_block(self, yaml_body: str) -> SlideData:
        """Your original regex parser, preserved as fallback."""
        slide: SlideData = {"layout": "title_and_content", "subtitle": "", "content": "", "left": "", "right": "", "image": ""}
        current_key: Optional[str] = None

        for line in yaml_body.splitlines():
            if not line.strip():
                if current_key and isinstance(slide.get(current_key), str):
                    slide[current_key] += "\n"
                continue

            is_indented = line.startswith("  ") or line.startswith("\t")
            match = re.match(r"^(\w+):\s*(.*)", line.strip())
            if match and not line.strip().startswith("-") and not is_indented:
                current_key = match.group(1).lower()
                value = match.group(2).strip()
                slide[current_key] = "" if value == "|" else value
            elif current_key:
                clean_line = line.rstrip()
                match_bullet = re.match(r"^(\s*)-\s*(.*)", clean_line)
                if match_bullet:
                    indent = match_bullet.group(1)
                    content = match_bullet.group(2)
                    clean_line = indent + content
                slide[current_key] = (slide[current_key] + "\n" + clean_line).strip() if slide[current_key] else clean_line

        for k, v in list(slide.items()):
            if isinstance(v, str):
                slide[k] = v.strip()
        return slide

    # ----------------------------
    # Text rendering (3 bullet levels)
    # ----------------------------

    def add_text_to_frame(self, text_frame, text: Union[str, List[str]]) -> None:
        """
        - Strips leading '-' bullet markers if present
        - Supports up to 3 levels (PowerPoint 0..2)
        """
        text_frame.clear()

        lines = text if isinstance(text, list) else str(text or "").splitlines()
        non_empty = [l for l in lines if str(l).strip()]
        if not non_empty:
            return

        # base indentation (avoid weird shifts)
        base_indent = min(len(l) - len(l.lstrip()) for l in non_empty)

        first = True
        for line in lines:
            if not str(line).strip():
                continue

            raw = str(line).rstrip("\n")
            stripped = raw.lstrip()
            indent = len(raw) - len(stripped)

            # strip markdown dash bullets
            stripped = re.sub(r"^-\s+", "", stripped)

            # 2 spaces => 1 level. cap to 3 levels total.
            level = max(0, min(((indent - base_indent) // 2), 2))

            p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            first = False

            p.text = stripped
            p.level = level

            if level == 0:
                p.font.size = Pt(22)
            elif level == 1:
                p.font.size = Pt(18)
            else:
                p.font.size = Pt(16)

    # ----------------------------
    # Proposal B: placeholder-based image placement
    # ----------------------------

    def add_picture_in_placeholder_box(self, slide, placeholder, img_path: str) -> None:
        img_full = _abs_from_plan(self.plan_path, img_path)
        if not os.path.exists(img_full):
            print(f"Warning: Image not found at {img_full}")
            return

        left, top, box_w, box_h = placeholder.left, placeholder.top, placeholder.width, placeholder.height

        # remove placeholder so it doesn't show "Click to add text"
        try:
            sp = placeholder.element
            sp.getparent().remove(sp)
        except Exception:
            pass

        w, h = _fit_image_into_box(img_full, box_w, box_h)
        pic_left = left + int((box_w - w) / 2)
        pic_top = top + int((box_h - h) / 2)

        slide.shapes.add_picture(img_full, pic_left, pic_top, width=w, height=h)

    # ----------------------------
    # Slide types
    # ----------------------------

    def add_title_slide(self, slide_data: SlideData) -> None:
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", "")

        subtitle = slide_data.get("subtitle", "")
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    def add_content_slide(self, slide_data: SlideData) -> None:
        # Title and Content
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", "")

        body = _find_placeholder(slide, idx=1)
        if body is not None:
            self.add_text_to_frame(body.text_frame, slide_data.get("content", ""))

        # image on the right half (template-independent fallback)
        img_path = str(slide_data.get("image", "") or "").strip()
        if img_path:
            img_full = _abs_from_plan(self.plan_path, img_path)
            if os.path.exists(img_full):
                slide_w = self.prs.slide_width
                slide_h = self.prs.slide_height

                # right column-ish box
                left = int(slide_w * 0.56)
                top = Inches(1.5)
                box_w = int(slide_w * 0.40)
                box_h = int(slide_h * 0.70)

                w, h = _fit_image_into_box(img_full, box_w, box_h)
                slide.shapes.add_picture(
                    img_full,
                    left + int((box_w - w) / 2),
                    top + int((box_h - h) / 2),
                    width=w,
                    height=h,
                )
            else:
                print(f"Warning: Image not found at {img_full}")

    def add_two_content_slide(self, slide_data: SlideData) -> None:
        # Two Content
        layout = self.prs.slide_layouts[3]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", "")

        left_placeholder = _find_placeholder(slide, idx=1)
        right_placeholder = _find_placeholder(slide, idx=2)

        if left_placeholder is not None:
            self.add_text_to_frame(left_placeholder.text_frame, slide_data.get("left", ""))

        right = slide_data.get("right", "")
        caption = ""

        # New: support YAML dict for right side
        if isinstance(right, dict):
            img = str(right.get("image") or right.get("img") or "").strip()
            caption = str(right.get("caption") or "").strip()
            right_text = right.get("text", "")
            if img and right_placeholder is not None:
                self.add_picture_in_placeholder_box(slide, right_placeholder, img)
            elif right_text and right_placeholder is not None:
                self.add_text_to_frame(right_placeholder.text_frame, right_text)
        else:
            right_text = str(right or "")
            if "image:" in right_text and right_placeholder is not None:
                img_match = re.search(r"image:\s*([^\n\r\s]+)", right_text)
                if img_match:
                    self.add_picture_in_placeholder_box(slide, right_placeholder, img_match.group(1).strip())
                cap_match = re.search(r"caption:\s*([^\n\r]+)", right_text)
                if cap_match:
                    caption = cap_match.group(1).strip()
            elif right_placeholder is not None:
                self.add_text_to_frame(right_placeholder.text_frame, right_text)

        if caption:
            footer = _find_placeholder(slide, placeholder_type=PP_PLACEHOLDER.FOOTER)
            if footer is not None:
                footer.text = caption
            else:
                box = slide.shapes.add_textbox(Inches(7.0), Inches(6.9), Inches(6.0), Inches(0.4))
                tf = box.text_frame
                tf.text = caption
                tf.paragraphs[0].font.size = Pt(12)

    # ----------------------------
    # Proposal C: smart layout upgrade
    # ----------------------------

    def generate(self) -> None:
        slides_data = self.parse_plan()

        for data in slides_data:
            layout_type = str(data.get("layout", "title_and_content") or "title_and_content").strip().lower()

            # Proposal C:
            # If the slide is title_and_content but has BOTH content and image,
            # render as Two Content (text left, image right) for better visual balance.
            if layout_type == "title_and_content":
                has_text = bool(str(data.get("content", "") or "").strip())
                has_img = bool(str(data.get("image", "") or "").strip())
                if has_text and has_img:
                    layout_type = "two_content"
                    data = dict(data)  # shallow copy
                    data["left"] = data.get("content", "")
                    cap = data.get("caption", "")
                    data["right"] = {"image": data.get("image", ""), "caption": cap}

            if layout_type == "title":
                self.add_title_slide(data)
            elif layout_type == "two_content":
                self.add_two_content_slide(data)
            else:
                # default to title_and_content
                self.add_content_slide(data)

        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        try:
            self.prs.save(self.output_path)
            print(f"Presentation saved to: {self.output_path}")
        except PermissionError:
            alt_path = self.output_path.replace(".pptx", "_alt.pptx")
            print(f"Warning: Could not save to {self.output_path} (file might be open). Saving to {alt_path} instead.")
            self.prs.save(alt_path)


if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.abspath(__file__))
    plan = os.path.join(base_dir, "presentation-slides-de.md")
    output = os.path.join(base_dir, "generated/SoftwareDevelopmentWithAI.pptx")
    files = os.path.join(base_dir, "files")
    template = os.path.join(base_dir, "template.pptx")

    os.makedirs(os.path.dirname(output), exist_ok=True)

    gen = PresentationGenerator(plan, output, files, template)
    gen.generate()
